from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from io import BytesIO
from typing import Dict, Any, List
import json
from app.services.llm_service import llm_service
from app.services.image_service import image_service

class PPTService:
    """PPT 生成服务"""

    def __init__(self):
        self.templates = self._load_templates()

    def _load_templates(self) -> Dict[str, Any]:
        """加载模板配置"""
        return {
            "modern-blue": {
                "name": "现代蓝色",
                "description": "简洁现代的蓝色主题",
                "thumbnail": "",
                "color_scheme": ["#4A90E2", "#50E3C2", "#B8E986", "#F5A623", "#E0E0E0"],
                "category": "商务",
                "background_color": "#FFFFFF",
                "title_color": "#4A90E2",
                "text_color": "#333333",
            },
            "elegant-purple": {
                "name": "优雅紫色",
                "description": "专业优雅的紫色主题",
                "thumbnail": "",
                "color_scheme": ["#9B59B6", "#3498DB", "#1ABC9C", "#F1C40F", "#E74C3C"],
                "category": "创意",
                "background_color": "#FFFFFF",
                "title_color": "#9B59B6",
                "text_color": "#333333",
            },
            "professional-gray": {
                "name": "专业灰色",
                "description": "稳重大气的灰色主题",
                "thumbnail": "",
                "color_scheme": ["#2C3E50", "#34495E", "#7F8C8D", "#95A5A6", "#BDC3C7"],
                "category": "商务",
                "background_color": "#FFFFFF",
                "title_color": "#2C3E50",
                "text_color": "#333333",
            },
            "vibrant-orange": {
                "name": "活力橙色",
                "description": "充满活力的橙色主题",
                "thumbnail": "",
                "color_scheme": ["#E67E22", "#F39C12", "#27AE60", "#2980B9", "#8E44AD"],
                "category": "创意",
                "background_color": "#FFFFFF",
                "title_color": "#E67E22",
                "text_color": "#333333",
            },
            "minimalist-green": {
                "name": "极简绿色",
                "description": "清新自然的绿色主题",
                "thumbnail": "",
                "color_scheme": ["#27AE60", "#2ECC71", "#16A085", "#1ABC9C", "#A3E4D7"],
                "category": "教育",
                "background_color": "#FFFFFF",
                "title_color": "#27AE60",
                "text_color": "#333333",
            },
            "dark-mode": {
                "name": "暗黑模式",
                "description": "科技感十足的暗黑主题",
                "thumbnail": "",
                "color_scheme": ["#1E1E1E", "#2D2D2D", "#3D3D3D", "#4D4D4D", "#5D5D5D"],
                "category": "科技",
                "background_color": "#1E1E1E",
                "title_color": "#FFFFFF",
                "text_color": "#E0E0E0",
            },
            "business-red": {
                "name": "商务红色",
                "description": "专业稳重的红色主题，适合正式场合",
                "thumbnail": "",
                "color_scheme": ["#C0392B", "#E74C3C", "#F1948A", "#E59866", "#D35400"],
                "category": "商务",
                "background_color": "#FFFFFF",
                "title_color": "#C0392B",
                "text_color": "#333333",
            },
            "business-gold": {
                "name": "商务金色",
                "description": "高端大气的金色主题，适合高端商务",
                "thumbnail": "",
                "color_scheme": ["#D4AC0D", "#F39C12", "#F1C40F", "#F7DC6F", "#F4D03F"],
                "category": "商务",
                "background_color": "#FFFFFF",
                "title_color": "#D4AC0D",
                "text_color": "#333333",
            },
            "education-blue": {
                "name": "教育蓝色",
                "description": "清新明亮的蓝色主题，适合教育培训",
                "thumbnail": "",
                "color_scheme": ["#2980B9", "#3498DB", "#5DADE2", "#85C1E9", "#AED6F1"],
                "category": "教育",
                "background_color": "#FFFFFF",
                "title_color": "#2980B9",
                "text_color": "#333333",
            },
            "education-yellow": {
                "name": "教育黄色",
                "description": "温暖活泼的黄色主题，适合儿童教育",
                "thumbnail": "",
                "color_scheme": ["#F39C12", "#F1C40F", "#F7DC6F", "#F4D03F", "#F8C471"],
                "category": "教育",
                "background_color": "#FFFFFF",
                "title_color": "#F39C12",
                "text_color": "#333333",
            },
            "tech-cyan": {
                "name": "科技青色",
                "description": "科技感十足的青色主题，适合科技展示",
                "thumbnail": "",
                "color_scheme": ["#16A085", "#1ABC9C", "#48C9B0", "#76D7C4", "#A3E4D7"],
                "category": "科技",
                "background_color": "#FFFFFF",
                "title_color": "#16A085",
                "text_color": "#333333",
            },
            "tech-purple": {
                "name": "科技紫色",
                "description": "未来感的紫色主题，适合科技产品",
                "thumbnail": "",
                "color_scheme": ["#8E44AD", "#9B59B6", "#AF7AC5", "#BB8FCE", "#D2B4DE"],
                "category": "科技",
                "background_color": "#FFFFFF",
                "title_color": "#8E44AD",
                "text_color": "#333333",
            },
            "creative-pink": {
                "name": "创意粉色",
                "description": "浪漫温馨的粉色主题，适合创意设计",
                "thumbnail": "",
                "color_scheme": ["#E91E63", "#EC407A", "#F06292", "#F48FB1", "#F8BBD0"],
                "category": "创意",
                "background_color": "#FFFFFF",
                "title_color": "#E91E63",
                "text_color": "#333333",
            },
            "creative-rainbow": {
                "name": "创意彩虹",
                "description": "多彩活力的彩虹主题，适合创意展示",
                "thumbnail": "",
                "color_scheme": ["#FF6B6B", "#4ECDC4", "#45B7D1", "#96CEB4", "#FFEAA7"],
                "category": "创意",
                "background_color": "#FFFFFF",
                "title_color": "#FF6B6B",
                "text_color": "#333333",
            },
            "finance-blue": {
                "name": "金融蓝色",
                "description": "专业稳重的金融主题，适合金融报告",
                "thumbnail": "",
                "color_scheme": ["#1A5276", "#2874A6", "#3498DB", "#5DADE2", "#85C1E9"],
                "category": "商务",
                "background_color": "#FFFFFF",
                "title_color": "#1A5276",
                "text_color": "#333333",
            },
            "medical-green": {
                "name": "医疗绿色",
                "description": "清新健康的医疗主题，适合医疗展示",
                "thumbnail": "",
                "color_scheme": ["#145A32", "#1E8449", "#27AE60", "#2ECC71", "#58D68D"],
                "category": "医疗",
                "background_color": "#FFFFFF",
                "title_color": "#145A32",
                "text_color": "#333333",
            },
            "legal-navy": {
                "name": "法律深蓝",
                "description": "严肃专业的法律主题，适合法律文档",
                "thumbnail": "",
                "color_scheme": ["#1B4F72", "#2874A6", "#3498DB", "#5499C7", "#85C1E9"],
                "category": "商务",
                "background_color": "#FFFFFF",
                "title_color": "#1B4F72",
                "text_color": "#333333",
            },
        }

    async def generate_ppt(self, outline_data: dict, template_id: str) -> dict:
        """生成 PPT"""
        template = self.templates.get(template_id, self.templates["modern-blue"])
        slides = []

        # 为每个章节生成幻灯片
        for section in outline_data.get("sections", []):
            # 生成幻灯片内容
            content = await llm_service.generate_slide_content(
                section["title"],
                section["content"]
            )

            # 搜索相关图片
            images = await image_service.search_images(section["title"], count=2)

            slides.append({
                "id": section["id"],
                "title": section["title"],
                "content": content,
                "layout": "title-content",
                "images": images,
            })

        return {
            "id": f"ppt_{outline_data.get('id', '')}",
            "title": outline_data.get("topic", "Untitled"),
            "slides": slides,
            "template": template,
            "created_at": outline_data.get("created_at", ""),
        }

    def create_pptx(self, ppt_data: dict) -> BytesIO:
        """创建 PPTX 文件"""
        prs = Presentation()
        template = ppt_data.get("template", {})

        # 设置幻灯片尺寸
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        for slide_data in ppt_data.get("slides", []):
            # 添加幻灯片
            slide_layout = prs.slide_layouts[1]  # 使用标题和内容布局
            slide = prs.slides.add_slide(slide_layout)

            # 设置背景色
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(template.get("background_color", "#FFFFFF")))

            # 添加标题
            title = slide.shapes.title
            title.text = slide_data["title"]
            title_frame = title.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(36)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(*self._hex_to_rgb(template.get("title_color", "#4A90E2")))
            title_para.alignment = PP_ALIGN.CENTER

            # 添加内容
            content_placeholder = slide.placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.clear()

            # 将内容分成段落
            paragraphs = slide_data["content"].split('\n')
            for i, para_text in enumerate(paragraphs):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = para_text.strip()
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(*self._hex_to_rgb(template.get("text_color", "#333333")))

            # 添加图片（如果有）
            if slide_data.get("images"):
                left = Inches(6)
                top = Inches(2)
                height = Inches(3)
                width = Inches(3.5)

                for img_url in slide_data["images"][:1]:  # 只添加第一张图片
                    try:
                        slide.shapes.add_picture(img_url, left, top, width=width, height=height)
                        break
                    except Exception as e:
                        print(f"Failed to add image: {e}")
                        continue

        # 保存到 BytesIO
        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer

    def _hex_to_rgb(self, hex_color: str) -> tuple:
        """将十六进制颜色转换为 RGB"""
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

    def get_templates(self) -> List[Dict]:
        """获取所有模板"""
        return [
            {
                "id": key,
                "name": value["name"],
                "description": value["description"],
                "thumbnail": value["thumbnail"],
                "colorScheme": value["color_scheme"],
                "category": value["category"],
            }
            for key, value in self.templates.items()
        ]

ppt_service = PPTService()